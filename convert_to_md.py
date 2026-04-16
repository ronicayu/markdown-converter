#!/usr/bin/env python3
"""Convert .doc, .docx, and .pdf files to Markdown.

PDFs use pymupdf4llm (better layout/bold/lists, image extraction).
Word docs use markitdown via LibreOffice (.doc -> .docx) then markitdown.
Images are processed with mlx-vlm (Qwen2.5-VL) for optional alt text and text extraction.

Usage:
    # Convert all files (DDS/, DRS/, Error Codes/, Message Specs V2.1/, MHX-MHub/, UserGuide/)
    .venv/bin/python3 convert_to_md.py

    # Convert a specific directory
    .venv/bin/python3 convert_to_md.py DDS/

    # Convert a single file
    .venv/bin/python3 convert_to_md.py "DDS/TN4.1 DDS - Liquor Tobacco/TN4.1 DDS - Liquor Tobacco (Rel1.0) UPD20130710.doc"

    # Dry run — show what would be converted
    .venv/bin/python3 convert_to_md.py --dry-run

    # Force reconvert even if .md already exists
    .venv/bin/python3 convert_to_md.py --force

    # Generate image alt text using mlx-vlm
    .venv/bin/python3 convert_to_md.py --describe-images

    # Extract text from images into sidecar .md files
    .venv/bin/python3 convert_to_md.py --extract-image-text
"""

import argparse
import hashlib
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import threading
import xml.etree.ElementTree as ET
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from zipfile import ZipFile

import pymupdf
import pymupdf4llm
from markitdown import MarkItDown


OLLAMA_ALT_PROMPT = (
    "Describe this screenshot from a technical user guide in one concise sentence, "
    "suitable for use as image alt text. Do not start with 'This image', 'The image' or 'A screenshot of'."
)

OLLAMA_EXTRACT_PROMPT = (
    "Extract all visible text from this screenshot exactly as it appears, preserving structure. "
    "Use markdown formatting: headings for labels, tables if tabular data is present, "
    "bullet points for lists. Output only the extracted text, no commentary."
)


_MIN_IMAGE_PX = 64  # skip images smaller than this in either dimension

# Ensures only one mlx-vlm generate() call runs at a time (Metal GPU is not re-entrant).
_mlx_serialize = threading.Lock()
# Set True after a timeout; prevents any further calls (the stuck thread still holds the lock).
_mlx_is_stuck = False


def describe_image(image_path: Path, model_and_processor: tuple, prompt: str = OLLAMA_ALT_PROMPT) -> str:
    """Call mlx-vlm vision model with the given prompt. Returns the response text.
    Returns empty string for images that are too small to be meaningful.

    Thread-safety: serializes all generate() calls via _mlx_serialize so that a
    timed-out daemon thread and a new call never run concurrently on Metal.  After
    a timeout _mlx_is_stuck is set and all subsequent calls return "" immediately.
    """
    global _mlx_is_stuck
    if _mlx_is_stuck:
        return ""

    # skip tiny images (icons, bullets, etc.)
    try:
        from PIL import Image as _Img
        w, h = _Img.open(image_path).size
        if w < _MIN_IMAGE_PX or h < _MIN_IMAGE_PX:
            return ""
    except Exception:
        pass

    try:
        from mlx_vlm import generate
        from mlx_vlm.prompt_utils import apply_chat_template

        model, processor, config = model_and_processor

        prompt_text = apply_chat_template(
            processor, config,
            prompt,
            num_images=1
        )

        result: list = [None]
        exc: list = [None]
        done = threading.Event()

        def _worker():
            with _mlx_serialize:  # blocks if a previous call is still running
                try:
                    result[0] = generate(
                        model, processor,
                        image=str(image_path),
                        prompt=prompt_text,
                        max_tokens=IMAGE_MAX_TOKENS,
                        verbose=False,
                    )
                except Exception as e:  # noqa: BLE001
                    exc[0] = e
            done.set()

        t = threading.Thread(target=_worker, daemon=True)
        t.start()
        if not done.wait(timeout=IMAGE_TIMEOUT):
            _mlx_is_stuck = True  # stuck thread still holds _mlx_serialize; skip all future images
            raise RuntimeError(
                f"mlx-vlm timed out ({IMAGE_TIMEOUT}s) for {image_path.name}; skipping remaining images"
            )
        if exc[0]:
            raise exc[0]

        response = result[0]
        # response is a GenerationResult object; extract text
        text = response.text if hasattr(response, 'text') else str(response)
        return text.strip()
    except Exception as e:
        raise RuntimeError(f"mlx-vlm request failed for {image_path.name}: {e}") from e


def detect_header_footer_margins(doc: pymupdf.Document) -> tuple[float, float]:
    """Detect header/footer heights by finding text blocks repeated across pages.

    Returns (top_margin, bottom_margin) in points. Returns (0, 0) for single-page docs.
    """
    if len(doc) < 2:
        return 0.0, 0.0

    # Collect text blocks per page: (y0, y1, text)
    per_page: list[tuple[float, list[tuple[float, float, str]]]] = []
    for page in doc:
        h = page.rect.height
        blocks = [
            (b[1], b[3], b[4].strip())
            for b in page.get_text("blocks")
            if b[6] == 0 and b[4].strip()
        ]
        per_page.append((h, blocks))

    # Count how many pages each text snippet appears on
    from collections import Counter
    text_count: Counter = Counter()
    text_pos: dict[str, tuple[float, float, float]] = {}  # key -> (y0, y1, page_height)
    for h, blocks in per_page:
        seen = set()
        for y0, y1, text in blocks:
            key = text.split("\n")[0][:80]  # first line only — avoids page-number suffix
            if key and key not in seen:
                text_count[key] += 1
                seen.add(key)
            if key and key not in text_pos:
                text_pos[key] = (y0, y1, h)

    # Text on 50%+ of pages and in top/bottom 20% of page → header/footer
    threshold = max(2, len(per_page) // 2)
    top_margin = 0.0
    bottom_margin = 0.0
    for key, count in text_count.items():
        if count < threshold:
            continue
        y0, y1, h = text_pos[key]
        if y0 < h * 0.20:
            top_margin = max(top_margin, y1 + 2)
        elif y0 > h * 0.80:
            bottom_margin = max(bottom_margin, h - y0 + 2)

    return top_margin, bottom_margin


def redact_header_footer(src: Path, dst: Path) -> None:
    """Redact header and footer zones from every page of a PDF (auto-detected)."""
    doc = pymupdf.open(str(src))
    top_pts, bottom_pts = detect_header_footer_margins(doc)
    if top_pts == 0 and bottom_pts == 0:
        doc.save(str(dst))
        doc.close()
        return
    for page in doc:
        w, h = page.rect.width, page.rect.height
        if top_pts:
            page.add_redact_annot(pymupdf.Rect(0, 0, w, top_pts))
        if bottom_pts:
            page.add_redact_annot(pymupdf.Rect(0, h - bottom_pts, w, h))
        page.apply_redactions(images=pymupdf.PDF_REDACT_IMAGE_REMOVE)
    doc.save(str(dst))
    doc.close()


def find_document_files(root: Path) -> list[Path]:
    """Find all .doc, .docx, .xlsx, .pptx, and .pdf files under root."""
    files = []
    for ext in ("*.doc", "*.docx", "*.xlsx", "*.pptx", "*.pdf"):
        files.extend(root.rglob(ext))
    return sorted(files)


def xlsx_to_md(src: Path) -> str:
    """Convert all sheets of an xlsx file to Markdown tables."""
    import openpyxl

    wb = openpyxl.load_workbook(src, data_only=True)
    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        # strip trailing all-None rows
        while rows and all(v is None for v in rows[-1]):
            rows.pop()
        if not rows:
            parts.append(f"## {sheet_name}\n\n*(empty)*")
            continue

        # determine column count from widest row
        ncols = max(len(r) for r in rows)

        def _cell(v) -> str:
            if v is None:
                return ""
            return str(v).replace("|", "\\|").replace("\n", " ").strip()

        header = rows[0]
        body = rows[1:]
        header_cells = [_cell(header[i] if i < len(header) else None) for i in range(ncols)]
        sep = ["---"] * ncols

        lines = [
            f"## {sheet_name}",
            "",
            "| " + " | ".join(header_cells) + " |",
            "| " + " | ".join(sep) + " |",
        ]
        for row in body:
            cells = [_cell(row[i] if i < len(row) else None) for i in range(ncols)]
            lines.append("| " + " | ".join(cells) + " |")
        parts.append("\n".join(lines))

    wb.close()
    return "\n\n".join(parts)


_CONFIG_PATH = Path(__file__).parent / "config.json"
_DEFAULT_PATHS = ["DRS", "Error Codes", "Message Specs V2.1", "MHX-MHub", "DDS", "UserGuide", "TN4.1 FE Vendor Guide"]


def _load_config_paths() -> list[str]:
    if _CONFIG_PATH.exists():
        with _CONFIG_PATH.open() as f:
            data = json.load(f)
        return data.get("paths", _DEFAULT_PATHS)
    return _DEFAULT_PATHS


SOFFICE = os.environ.get("SOFFICE_PATH", "/Applications/LibreOffice.app/Contents/MacOS/soffice")


SOFFICE_TIMEOUT = int(os.environ.get("SOFFICE_TIMEOUT", "120"))

IMAGE_TIMEOUT = int(os.environ.get("IMAGE_TIMEOUT", "120"))
IMAGE_MAX_TOKENS = int(os.environ.get("IMAGE_MAX_TOKENS", "4096"))


def doc_to_docx(doc_path: Path, tmp_dir: Path) -> Path:
    """Convert .doc to .docx using LibreOffice headless."""
    try:
        result = subprocess.run(
            [SOFFICE, "--headless", "--convert-to", "docx", "--outdir", str(tmp_dir), str(doc_path)],
            capture_output=True,
            text=True,
            timeout=SOFFICE_TIMEOUT,
        )
    except subprocess.TimeoutExpired:
        raise RuntimeError(f"LibreOffice timed out ({SOFFICE_TIMEOUT}s) for {doc_path}")
    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice failed for {doc_path}: {result.stderr.strip()}")
    out_path = tmp_dir / (doc_path.stem + ".docx")
    if not out_path.exists():
        raise RuntimeError(f"LibreOffice produced no output for {doc_path}")
    return out_path


PLACEHOLDER_RE = re.compile(r"!\[([^\]]*)\]\(data:image/[^)]*\.\.\.\)")


def get_header_footer_media(z: ZipFile) -> set[str]:
    """Return set of word/media/... paths referenced by header or footer parts."""
    hf_media: set[str] = set()
    hf_parts = [n for n in z.namelist() if re.match(r"word/(header|footer)\d*\.xml$", n)]
    for part in hf_parts:
        rels_path = part.replace("word/", "word/_rels/") + ".rels"
        if rels_path not in z.namelist():
            continue
        try:
            tree = ET.fromstring(z.read(rels_path))
            part_dir = part.rsplit("/", 1)[0] + "/"  # e.g. "word/"
            for rel in tree:
                target = rel.get("Target", "")
                if "media/" in target:
                    # resolve relative to part's directory
                    hf_media.add(part_dir + target.lstrip("/").replace("../", ""))
        except ET.ParseError:
            pass
    return hf_media


# Formats Ollama vision models cannot handle — will be converted to PNG
_VECTOR_EXTS = {".wmf", ".emf", ".svg"}
# Formats safe to send to Ollama
_RASTER_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}


def _to_png(data: bytes, src_ext: str) -> bytes | None:
    """Convert vector image bytes to PNG using pymupdf. Returns None on failure."""
    try:
        doc = pymupdf.open(stream=data, filetype=src_ext.lstrip("."))
        pix = doc[0].get_pixmap(dpi=150)
        return pix.tobytes("png")
    except Exception:
        return None


def extract_images_from_docx(docx_path: Path, img_dir: Path, md_dir: Path) -> list[str]:
    """Extract body images from a docx into img_dir, skipping header/footer images.
    Vector formats (wmf/emf/svg) are converted to PNG.
    Returns paths relative to md_dir (e.g. images/doc-slug/001.png)."""
    image_paths = []
    with ZipFile(docx_path) as z:
        hf_media = get_header_footer_media(z)
        media = sorted(n for n in z.namelist() if n.startswith("word/media/") and n not in hf_media)
        if not media:
            return image_paths
        img_dir.mkdir(parents=True, exist_ok=True)
        for i, name in enumerate(media, 1):
            ext = Path(name).suffix.lower()
            data = z.read(name)
            if ext in _VECTOR_EXTS:
                png_data = _to_png(data, ext)
                if png_data:
                    filename = f"{i:03d}.png"
                    (img_dir / filename).write_bytes(png_data)
                else:
                    continue  # skip unconvertible vector images
            else:
                filename = f"{i:03d}{ext}"
                (img_dir / filename).write_bytes(data)
            image_paths.append(str((img_dir / filename).relative_to(md_dir)))
    return image_paths


_CONTENT_TYPE_EXT = {
    "image/png": ".png",
    "image/jpeg": ".jpg",
    "image/gif": ".gif",
    "image/bmp": ".bmp",
    "image/webp": ".webp",
    "image/x-wmf": ".wmf",
    "image/x-emf": ".emf",
    "image/svg+xml": ".svg",
}


def extract_images_from_pptx(pptx_path: Path, img_dir: Path, md_dir: Path) -> dict[str, str]:
    """Extract images from a pptx using python-pptx, keyed by markitdown's synthesised filename.

    markitdown generates image refs as ``re.sub(r"\\W", "", shape.name) + ".jpg"``
    regardless of the real image format.  We extract the actual blob with its real
    extension and return {markitdown_ref → real_relative_path} so the caller can
    replace the broken references in the markdown.
    """
    import pptx as _pptx
    import pptx.enum.shapes as _shapes

    mapping: dict[str, str] = {}
    img_dir.mkdir(parents=True, exist_ok=True)
    prs = _pptx.Presentation(str(pptx_path))

    def _extract(shape):
        is_pic = (
            shape.shape_type == _shapes.MSO_SHAPE_TYPE.PICTURE
            or (shape.shape_type == _shapes.MSO_SHAPE_TYPE.PLACEHOLDER and hasattr(shape, "image"))
        )
        if is_pic:
            ref_name = re.sub(r"\W", "", shape.name) + ".jpg"  # matches markitdown exactly
            if ref_name not in mapping:
                ct = (shape.image.content_type or "image/png").lower().split(";")[0].strip()
                real_ext = _CONTENT_TYPE_EXT.get(ct, ".png")
                stem = re.sub(r"\W", "", shape.name)
                blob = shape.image.blob
                if real_ext in _VECTOR_EXTS:
                    png = _to_png(blob, real_ext)
                    if png:
                        fname = stem + ".png"
                        (img_dir / fname).write_bytes(png)
                        mapping[ref_name] = str((img_dir / fname).relative_to(md_dir))
                else:
                    fname = stem + real_ext
                    (img_dir / fname).write_bytes(blob)
                    mapping[ref_name] = str((img_dir / fname).relative_to(md_dir))

        if shape.shape_type == _shapes.MSO_SHAPE_TYPE.GROUP:
            for sub in shape.shapes:
                _extract(sub)

    for slide in prs.slides:
        for shape in slide.shapes:
            _extract(shape)

    return mapping


def replace_image_placeholders(md_text: str, image_paths: list[str]) -> str:
    """Replace data:image/...base64... placeholders with extracted file paths."""
    idx = 0

    def _replace(m: re.Match) -> str:
        nonlocal idx
        alt = m.group(1)
        if idx < len(image_paths):
            path = image_paths[idx]
            idx += 1
            return f"![{alt}]({path})"
        return m.group(0)  # no more images, keep placeholder

    return PLACEHOLDER_RE.sub(_replace, md_text)


def _apply_vision_model(final_img: Path, img_desc_fn, img_extract_fn) -> str:
    """Generate alt text and/or write sidecar .md for an image. Returns alt text."""
    alt = ""
    if img_desc_fn is not None:
        try:
            alt = img_desc_fn(final_img)
        except Exception as e:
            print(f"\n  [warn] image description failed ({final_img.name}): {e}", file=sys.stderr)
    if img_extract_fn is not None:
        sidecar = final_img.with_suffix(".md")
        if not sidecar.exists():
            try:
                extracted = img_extract_fn(final_img)
                if extracted:
                    sidecar.write_text(extracted, encoding="utf-8")
            except Exception as e:
                print(f"\n  [warn] image text extraction failed ({final_img.name}): {e}", file=sys.stderr)
    return alt


def convert_file(md_converter: MarkItDown, src: Path, dest: Path, tmp_dir: Path, img_desc_fn=None, img_extract_fn=None, workers: int = 1) -> dict:
    """Convert a single Word or PDF file to Markdown. Returns status dict."""
    try:
        suffix = src.suffix.lower()
        dest.parent.mkdir(parents=True, exist_ok=True)

        if suffix == ".pdf":
            doc_slug = normalize_stem(src.stem)
            img_dir = dest.parent / "images" / doc_slug
            shutil.rmtree(img_dir, ignore_errors=True)
            img_dir.mkdir(parents=True, exist_ok=True)
            clean_pdf = tmp_dir / src.name
            redact_header_footer(src, clean_pdf)

            # write images to a per-file temp subdir to avoid name collisions
            file_hash = hashlib.md5(src.stem.encode()).hexdigest()[:6]
            tmp_img_dir = (tmp_dir / file_hash).resolve()
            tmp_img_dir.mkdir(exist_ok=True)
            md_text = pymupdf4llm.to_markdown(
                str(clean_pdf),
                write_images=True,
                image_path=str(tmp_img_dir),
                image_format="png",
            )

            # rename images to short names in order of appearance in markdown
            seen: dict[str, str] = {}  # abs_path -> images/doc_slug/short_name
            img_count = 0
            for m in re.finditer(r'\]\(([^)]+\.png)\)', md_text):
                abs_path = m.group(1)
                if abs_path in seen:
                    continue
                img_count += 1
                new_name = f"{img_count:03d}.png"
                src_img = Path(abs_path)
                if src_img.exists():
                    src_img.rename(img_dir / new_name)
                seen[abs_path] = f"images/{doc_slug}/{new_name}"

            # replace paths and optionally add alt text / extract text via vision model
            def _process_pdf_img(item: tuple[str, str]) -> tuple[str, str, str | None]:
                abs_path, rel_path = item
                final_img = img_dir / Path(rel_path).name
                if final_img.suffix.lower() not in _RASTER_EXTS:
                    return abs_path, rel_path, None  # None = non-raster, just repath
                return abs_path, rel_path, _apply_vision_model(final_img, img_desc_fn, img_extract_fn)

            _show_img_progress = (img_desc_fn is not None or img_extract_fn is not None) and bool(seen)
            with ThreadPoolExecutor(max_workers=workers) as img_pool:
                futures = {img_pool.submit(_process_pdf_img, item): item for item in seen.items()}
                n_done = 0
                for fut in as_completed(futures):
                    n_done += 1
                    if _show_img_progress:
                        print(f" [{n_done}/{len(futures)}]", end="", flush=True)
                    abs_path, rel_path, alt = fut.result()
                    if alt is None:
                        md_text = md_text.replace(f"]({abs_path})", f"]({rel_path})")
                    else:
                        md_text = md_text.replace(f"]({abs_path})", f"{alt}]({rel_path})")
        elif suffix == ".xlsx":
            md_text = xlsx_to_md(src)
            img_count = 0
        elif suffix == ".pptx":
            result = md_converter.convert(str(src))
            img_dir = dest.parent / "images" / normalize_stem(src.stem)
            shutil.rmtree(img_dir, ignore_errors=True)
            img_mapping = extract_images_from_pptx(src, img_dir, dest.parent)
            md_text = result.text_content
            # markitdown references images by their original basename — repath to extracted location
            for orig_name, new_path in img_mapping.items():
                md_text = md_text.replace(f"]({orig_name})", f"]({new_path})")
            image_paths = list(img_mapping.values())
            img_count = len(image_paths)

            def _process_pptx_img(rel_path: str) -> tuple[str, str]:
                final_img = dest.parent / rel_path
                if not final_img.exists() or final_img.suffix.lower() not in _RASTER_EXTS:
                    return rel_path, ""
                return rel_path, _apply_vision_model(final_img, img_desc_fn, img_extract_fn)

            _show_img_progress = (img_desc_fn is not None or img_extract_fn is not None) and bool(image_paths)
            with ThreadPoolExecutor(max_workers=workers) as img_pool:
                futures = {img_pool.submit(_process_pptx_img, rp): rp for rp in image_paths}
                n_done = 0
                for fut in as_completed(futures):
                    n_done += 1
                    if _show_img_progress:
                        print(f" [{n_done}/{len(futures)}]", end="", flush=True)
                    rel_path, alt = fut.result()
                    if alt:
                        md_text = md_text.replace(f"]({rel_path})", f"{alt}]({rel_path})")
        else:
            if suffix == ".doc":
                docx_path = doc_to_docx(src, tmp_dir)
            else:
                docx_path = src
            result = md_converter.convert(str(docx_path))
            img_dir = dest.parent / "images" / normalize_stem(src.stem)
            shutil.rmtree(img_dir, ignore_errors=True)
            image_paths = extract_images_from_docx(docx_path, img_dir, dest.parent)
            md_text = replace_image_placeholders(result.text_content, image_paths)
            img_count = len(image_paths)
            if suffix == ".doc":
                docx_path.unlink(missing_ok=True)

            # apply vision model to docx images
            def _process_docx_img(rel_path: str) -> tuple[str, str]:
                final_img = dest.parent / rel_path
                if not final_img.exists() or final_img.suffix.lower() not in _RASTER_EXTS:
                    return rel_path, ""
                return rel_path, _apply_vision_model(final_img, img_desc_fn, img_extract_fn)

            _show_img_progress = (img_desc_fn is not None or img_extract_fn is not None) and bool(image_paths)
            with ThreadPoolExecutor(max_workers=workers) as img_pool:
                futures = {img_pool.submit(_process_docx_img, rp): rp for rp in image_paths}
                n_done = 0
                for fut in as_completed(futures):
                    n_done += 1
                    if _show_img_progress:
                        print(f" [{n_done}/{len(futures)}]", end="", flush=True)
                    rel_path, alt = fut.result()
                    if alt:
                        md_text = md_text.replace(f"]({rel_path})", f"{alt}]({rel_path})")

        dest.write_text(md_text, encoding="utf-8")
        return {"status": "ok", "src": src, "dest": dest, "size": len(md_text), "images": img_count}

    except Exception as e:
        return {"status": "error", "src": src, "dest": dest, "error": str(e)}


def normalize_stem(stem: str) -> str:
    """Normalize a filename stem to a URL-safe, shell-friendly slug."""
    s = re.sub(r'[()[\]{}]', '', stem)  # remove brackets
    s = s.replace('&', '-')             # & → -
    s = re.sub(r'[ _]', '-', s)         # spaces/underscores → -
    s = re.sub(r'\.', '-', s)           # dots → -
    s = re.sub(r'-+', '-', s)           # collapse runs
    return s.strip('-').lower()


def md_output_path(src: Path, output_dir: Path | None = None) -> Path:
    """Determine .md output path, optionally rooted under output_dir."""
    rel = src.with_name(normalize_stem(src.stem) + ".md")
    if output_dir is not None:
        return output_dir / rel
    return rel


def main():
    parser = argparse.ArgumentParser(description="Convert Word documents to Markdown")
    parser.add_argument(
        "paths",
        nargs="*",
        default=None,
        help="Files or directories to convert (default: paths from convert_scripts/config.json)",
    )
    parser.add_argument("--dry-run", action="store_true", help="Show what would be converted")
    parser.add_argument("--force", action="store_true", help="Reconvert even if .md exists")
    parser.add_argument("--output-dir", metavar="DIR", default="out", help="Root output directory (mirrors source structure, default: out/)")
    parser.add_argument("--describe-images", action="store_true", help="Use mlx-vlm vision model to generate image alt text")
    parser.add_argument("--extract-image-text", action="store_true", help="Use mlx-vlm vision model to extract text from each image into a sidecar .md file")
    parser.add_argument("--workers", type=int, default=1, metavar="N", help="Parallel image-processing workers per file (default: 1)")
    parser.add_argument("--images-only", action="store_true", help="Skip document conversion; only run vision model on existing images")
    args = parser.parse_args()
    if not args.paths:
        args.paths = _load_config_paths()

    output_dir = Path(args.output_dir)

    img_desc_fn = None
    img_extract_fn = None
    model_and_processor = None
    if args.describe_images or args.extract_image_text:
        from functools import partial
        try:
            print("Loading mlx-vlm Qwen2.5-VL-7B-Instruct-4bit ...", flush=True)
            from mlx_vlm import load
            from mlx_vlm.utils import load_config
            model, processor = load("mlx-community/Qwen2.5-VL-7B-Instruct-4bit")
            config = load_config("mlx-community/Qwen2.5-VL-7B-Instruct-4bit")
            model_and_processor = (model, processor, config)
            print("Model loaded successfully")
        except Exception as e:
            print(f"Error loading mlx-vlm: {e}", file=sys.stderr)
            sys.exit(1)
        if args.describe_images:
            img_desc_fn = partial(describe_image, model_and_processor=model_and_processor)
        if args.extract_image_text:
            img_extract_fn = partial(describe_image, model_and_processor=model_and_processor, prompt=OLLAMA_EXTRACT_PROMPT)

    # Collect files
    files: list[Path] = []
    for p in args.paths:
        path = Path(p)
        if path.is_file():
            files.append(path)
        elif path.is_dir():
            files.extend(find_document_files(path))
        else:
            print(f"Warning: {p} not found, skipping", file=sys.stderr)

    if not files:
        print("No Word files found.")
        return

    # Filter already-converted unless --force
    to_convert = []
    already_done = []
    skipped = 0
    for f in files:
        dest = md_output_path(f, output_dir)
        if dest.exists() and not args.force:
            skipped += 1
            already_done.append((f, dest))
        else:
            to_convert.append((f, dest))

    if args.images_only:
        already_done.extend(to_convert)
        to_convert = []

    print(f"Found {len(files)} Word files, {len(to_convert)} to convert, {skipped} already done")

    if args.dry_run:
        for src, dest in to_convert:
            print(f"  {src} -> {dest}")
        return

    if not to_convert and not args.images_only:
        print("Nothing to do.")
        return

    md_converter = MarkItDown() if to_convert else None
    ok_count = 0
    err_count = 0
    errors = []

    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp_path = Path(tmp_dir)
        for i, (src, dest) in enumerate(to_convert, 1):
            print(f"[{i}/{len(to_convert)}] {src.name} ... ", end="", flush=True)
            result = convert_file(md_converter, src, dest, tmp_path, img_desc_fn, img_extract_fn, args.workers)
            if result["status"] == "ok":
                ok_count += 1
                imgs = result.get("images", 0)
                img_info = f", {imgs} images" if imgs else ""
                print(f"OK ({result['size']:,} chars{img_info})")
            else:
                err_count += 1
                errors.append(result)
                print(f"FAILED: {result['error']}")

    print(f"\nDone: {ok_count} converted, {err_count} failed, {skipped} skipped")
    if errors:
        print("\nFailed files:")
        for e in errors:
            print(f"  {e['src']}: {e['error']}")

    # For already-converted docs, process any missing image sidecars
    if img_extract_fn is not None and already_done:
        pending = [
            (src, dest) for src, dest in already_done
            if any(
                True for img in (dest.parent / "images" / normalize_stem(src.stem)).glob("*")
                if img.suffix.lower() in _RASTER_EXTS and not img.with_suffix(".md").exists()
            )
        ]
        if pending:
            print(f"\nProcessing image sidecars for {len(pending)} already-converted files ...")
            for src, dest in pending:
                img_dir = dest.parent / "images" / normalize_stem(src.stem)
                imgs = sorted(
                    i for i in img_dir.glob("*")
                    if i.suffix.lower() in _RASTER_EXTS and not i.with_suffix(".md").exists()
                )
                if not imgs:
                    continue
                prefix = f"  {src.name} ({len(imgs)} images)"
                print(f"{prefix} ... ", end="", flush=True)

                def _run_sidecar(img: Path) -> bool:
                    print(f"\r{prefix} processing {img.name:<12}", end="", flush=True)
                    wrote = False
                    try:
                        extracted = img_extract_fn(img)
                        if extracted:
                            img.with_suffix(".md").write_text(extracted, encoding="utf-8")
                            wrote = True
                    except Exception as e:
                        print(f"\n    [warn] {img.name}: {e}", file=sys.stderr)
                    return wrote

                with ThreadPoolExecutor(max_workers=args.workers) as img_pool:
                    fut_to_img = {img_pool.submit(_run_sidecar, img): img for img in imgs}
                    n_done = 0
                    done = 0
                    for fut in as_completed(fut_to_img):
                        n_done += 1
                        img_name = fut_to_img[fut].name
                        print(f"\r{prefix} [{n_done}/{len(fut_to_img)}] {img_name:<12}", end="", flush=True)
                        if fut.result():
                            done += 1
                print(f"\r{prefix}: {done} new sidecars{' ' * 20}")


if __name__ == "__main__":
    main()
